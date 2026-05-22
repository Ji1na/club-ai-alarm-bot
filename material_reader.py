import io
import httpx
from bs4 import BeautifulSoup
from pptx import Presentation
from pypdf import PdfReader


def _clean_text(text: str, limit: int = 15000) -> str:
    lines = [line.strip() for line in text.splitlines() if line.strip()]
    return "\n".join(lines)[:limit]


def read_material(url: str | None) -> str | None:
    """
    URL에서 텍스트 내용을 읽어옵니다.
    - 일반 웹페이지: HTML에서 텍스트 추출
    - 텍스트 파일: 그대로 반환
    - URL 없으면 None 반환

    이 함수는 '참고자료 URL'을 받을 때 쓰는 보조 기능입니다.
    """
    if not url:
        return None

    try:
        with httpx.Client(timeout=10, follow_redirects=True) as client:
            r = client.get(url)
            r.raise_for_status()
            content_type = r.headers.get("content-type", "")

            if "text/plain" in content_type:
                return _clean_text(r.text)

            if "text/html" in content_type:
                soup = BeautifulSoup(r.text, "html.parser")

                for tag in soup(["script", "style", "nav", "footer", "header"]):
                    tag.decompose()

                text = soup.get_text(separator="\n", strip=True)
                return _clean_text(text)

            return _clean_text(r.text)

    except Exception as e:
        print(f"[material_reader] URL 읽기 실패: {e}")
        return None


async def extract_text_from_discord_attachment(attachment) -> str:
    """
    Discord slash command에서 받은 첨부파일을 텍스트로 변환합니다.

    지원 형식:
    - .txt
    - .md
    - .pptx
    - .pdf
    """
    data = await attachment.read()
    filename = attachment.filename.lower()

    if filename.endswith(".txt") or filename.endswith(".md"):
        return _clean_text(data.decode("utf-8", errors="ignore"))

    if filename.endswith(".pptx"):
        return extract_text_from_pptx(data)

    if filename.endswith(".pdf"):
        return extract_text_from_pdf(data)

    raise ValueError("지원하지 않는 파일 형식입니다. pptx, pdf, txt, md만 지원합니다.")


def extract_text_from_pptx(data: bytes) -> str:
    prs = Presentation(io.BytesIO(data))
    chunks = []

    for slide_index, slide in enumerate(prs.slides, start=1):
        texts = []

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())

        if texts:
            chunks.append(f"[Slide {slide_index}]\n" + "\n".join(texts))

    return _clean_text("\n\n".join(chunks))


def extract_text_from_pdf(data: bytes) -> str:
    reader = PdfReader(io.BytesIO(data))
    chunks = []

    for page_index, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""

        if text.strip():
            chunks.append(f"[Page {page_index}]\n{text.strip()}")

    return _clean_text("\n\n".join(chunks))